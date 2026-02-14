"""
Document Parsers - Extract text from various file formats.
Supports: TXT, JSON, PDF, XML, DOC/DOCX, ODT/ODS/ODP, XLS/XLSX, PPTX, Images
"""
import io
import json
import logging
from pathlib import Path
from typing import Optional
logger = logging.getLogger("DocumentParsers")

class DocumentParser:
    """Universal document parser supporting multiple formats."""
    @staticmethod
    def parse(filename: str, content: bytes) -> str:
        """Parse document content based on file extension."""
        ext = Path(filename).suffix.lower()
        parsers = {
            '.txt': DocumentParser._parse_text,
            '.json': DocumentParser._parse_json,
            '.pdf': DocumentParser._parse_pdf,
            '.xml': DocumentParser._parse_xml,
            '.doc': DocumentParser._parse_doc,
            '.docx': DocumentParser._parse_docx,
            '.odt': DocumentParser._parse_odt,
            '.ods': DocumentParser._parse_ods,
            '.odp': DocumentParser._parse_odp,
            '.xls': DocumentParser._parse_xls,
            '.xlsx': DocumentParser._parse_xlsx,
            '.pptx': DocumentParser._parse_pptx,
            '.png': DocumentParser._parse_image,
            '.jpg': DocumentParser._parse_image,
            '.jpeg': DocumentParser._parse_image,
            '.gif': DocumentParser._parse_image,
            '.webp': DocumentParser._parse_image,
        }
        parser_fn = parsers.get(ext)
        if not parser_fn:
            raise ValueError(f"Unsupported file format: {ext}")
        try:
            return parser_fn(content, filename)
        except Exception as e:
            logger.error(f"Failed to parse {filename}: {e}")
            raise
    @staticmethod
    def _parse_text(content: bytes, filename: str = "") -> str:
        return content.decode('utf-8', errors='replace')
    @staticmethod
    def _parse_json(content: bytes, filename: str = "") -> str:
        data = json.loads(content.decode('utf-8'))
        return json.dumps(data, indent=2, ensure_ascii=False)
    @staticmethod
    def _parse_pdf(content: bytes, filename: str = "") -> str:
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(content))
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return '\n\n'.join(text_parts)
    @staticmethod
    def _parse_xml(content: bytes, filename: str = "") -> str:
        from lxml import etree
        tree = etree.fromstring(content)
        return etree.tostring(tree, pretty_print=True, encoding='unicode')
    @staticmethod
    def _parse_doc(content: bytes, filename: str = "") -> str:
        """Basic .doc parsing - try as text or use antiword if available."""
        try:
            import subprocess
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.doc', delete=True) as tmp:
                tmp.write(content)
                tmp.flush()
                result = subprocess.run(
                    ['antiword', tmp.name], capture_output=True, text=True, timeout=30
                )
                if result.returncode == 0:
                    return result.stdout
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass
        # Fallback: try as plain text
        return content.decode('utf-8', errors='replace')
    @staticmethod
    def _parse_docx(content: bytes, filename: str = "") -> str:
        from docx import Document
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                paragraphs.append(' | '.join(cells))
        return '\n\n'.join(paragraphs)
    @staticmethod
    def _parse_odt(content: bytes, filename: str = "") -> str:
        from odf.opendocument import load
        from odf.text import P
        from odf import teletype
        doc = load(io.BytesIO(content))
        paragraphs = []
        for p in doc.getElementsByType(P):
            text = teletype.extractText(p)
            if text.strip():
                paragraphs.append(text)
        return '\n\n'.join(paragraphs)
    @staticmethod
    def _parse_ods(content: bytes, filename: str = "") -> str:
        from odf.opendocument import load
        from odf.table import Table, TableRow, TableCell
        from odf import teletype
        doc = load(io.BytesIO(content))
        lines = []
        for table in doc.getElementsByType(Table):
            table_name = table.getAttribute('name') or 'Sheet'
            lines.append(f"--- {table_name} ---")
            for row in table.getElementsByType(TableRow):
                cells = []
                for cell in row.getElementsByType(TableCell):
                    cells.append(teletype.extractText(cell))
                lines.append(' | '.join(cells))
        return '\n'.join(lines)
    @staticmethod
    def _parse_odp(content: bytes, filename: str = "") -> str:
        from odf.opendocument import load
        from odf.text import P
        from odf import teletype
        doc = load(io.BytesIO(content))
        texts = []
        for p in doc.getElementsByType(P):
            text = teletype.extractText(p)
            if text.strip():
                texts.append(text)
        return '\n\n'.join(texts)
    @staticmethod
    def _parse_xls(content: bytes, filename: str = "") -> str:
        """Parse legacy .xls files."""
        try:
            import openpyxl
            # Try loading as xlsx first (some .xls are actually xlsx)
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True)
            return DocumentParser._extract_openpyxl(wb)
        except Exception:
            # Fallback: try as text
            return content.decode('utf-8', errors='replace')
    @staticmethod
    def _parse_xlsx(content: bytes, filename: str = "") -> str:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True)
        return DocumentParser._extract_openpyxl(wb)
    @staticmethod
    def _extract_openpyxl(wb) -> str:
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"--- {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else '' for c in row]
                lines.append(' | '.join(cells))
        return '\n'.join(lines)
    @staticmethod
    def _parse_pptx(content: bytes, filename: str = "") -> str:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)
            slides_text.append('\n'.join(slide_parts))
        return '\n\n'.join(slides_text)
    @staticmethod
    def _parse_image(content: bytes, filename: str = "") -> str:
        """Extract metadata and OCR text from images."""
        from PIL import Image
        img = Image.open(io.BytesIO(content))
        info_parts = [
            f"[Image: {filename}]",
            f"Size: {img.size[0]}x{img.size[1]}",
            f"Format: {img.format}",
            f"Mode: {img.mode}"
        ]
        # Try OCR with pytesseract if available
        try:
            import pytesseract
            text = pytesseract.image_to_string(img)
            if text.strip():
                info_parts.append(f"OCR Text: {text.strip()}")
        except (ImportError, Exception):
            info_parts.append("(OCR not available)")
        return '\n'.join(info_parts)
